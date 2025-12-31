"""
문서 파서 - PDF, DOCX, PPTX, XLSX 파싱
"""

import os
from pathlib import Path
from typing import Dict, List, Optional
import logging

# PDF 파싱
import fitz  # PyMuPDF

# Office 문서 파싱
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

# OCR (optional)
try:
    from paddleocr import PaddleOCR
    PADDLE_AVAILABLE = True
except ImportError:
    PADDLE_AVAILABLE = False
    logger = logging.getLogger(__name__)
    logging.warning("PaddleOCR not available. Image parsing will be disabled.")

from PIL import Image

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class DocumentParser:
    """다양한 문서 포맷을 파싱하는 통합 파서"""

    def __init__(self, ocr_config: Optional[Dict] = None):
        """
        Args:
            ocr_config: OCR 설정 딕셔너리
        """
        self.ocr_config = ocr_config or {
            "lang": "korean",
            "use_angle_cls": True,
            "use_gpu": True,
            "show_log": False
        }

        # OCR 모델 초기화 (lazy loading)
        self._ocr_model = None

    @property
    def ocr_model(self):
        """OCR 모델 lazy loading"""
        if not PADDLE_AVAILABLE:
            raise ImportError("PaddleOCR is not installed. Install with: pip install paddlepaddle paddleocr")
        if self._ocr_model is None:
            logger.info("PaddleOCR 모델 로딩 중...")
            self._ocr_model = PaddleOCR(**self.ocr_config)
        return self._ocr_model

    def parse_file(self, file_path: str) -> Dict[str, any]:
        """
        파일을 파싱하여 텍스트 추출

        Args:
            file_path: 파일 경로

        Returns:
            파싱 결과 딕셔너리
            {
                'text': str,
                'metadata': dict,
                'success': bool,
                'error': str (if failed)
            }
        """
        file_path = Path(file_path)

        if not file_path.exists():
            return {
                'text': '',
                'metadata': {},
                'success': False,
                'error': f'File not found: {file_path}'
            }

        file_ext = file_path.suffix.lower()

        try:
            if file_ext == '.pdf':
                result = self.parse_pdf(str(file_path))
            elif file_ext == '.docx':
                result = self.parse_docx(str(file_path))
            elif file_ext == '.pptx':
                result = self.parse_pptx(str(file_path))
            elif file_ext == '.xlsx':
                result = self.parse_xlsx(str(file_path))
            elif file_ext in ['.jpg', '.jpeg', '.png']:
                result = self.parse_image(str(file_path))
            else:
                return {
                    'text': '',
                    'metadata': {},
                    'success': False,
                    'error': f'Unsupported file format: {file_ext}'
                }

            # 메타데이터 추가
            result['metadata']['file_name'] = file_path.name
            result['metadata']['file_path'] = str(file_path.absolute())
            result['metadata']['file_type'] = file_ext[1:]  # . 제거
            result['metadata']['file_size'] = file_path.stat().st_size

            return result

        except Exception as e:
            logger.error(f"Error parsing {file_path}: {e}")
            return {
                'text': '',
                'metadata': {},
                'success': False,
                'error': str(e)
            }

    def parse_pdf(self, file_path: str) -> Dict:
        """PDF 파일 파싱 (PyMuPDF 사용)"""
        try:
            doc = fitz.open(file_path)
            text_content = []
            metadata = {
                'page_count': len(doc),
                'has_images': False
            }

            for page_num in range(len(doc)):
                page = doc[page_num]
                text = page.get_text()
                text_content.append(text)

                # 이미지 포함 여부 확인
                if page.get_images():
                    metadata['has_images'] = True

            doc.close()

            full_text = '\n\n'.join(text_content)

            return {
                'text': full_text,
                'metadata': metadata,
                'success': True
            }

        except Exception as e:
            logger.error(f"PDF parsing error: {e}")
            return {
                'text': '',
                'metadata': {},
                'success': False,
                'error': str(e)
            }

    def parse_docx(self, file_path: str) -> Dict:
        """Word 문서 파싱"""
        try:
            doc = Document(file_path)
            text_content = []

            # 단락 추출
            for para in doc.paragraphs:
                if para.text.strip():
                    text_content.append(para.text)

            # 표 추출
            for table in doc.tables:
                for row in table.rows:
                    row_text = '\t'.join([cell.text for cell in row.cells])
                    text_content.append(row_text)

            metadata = {
                'paragraph_count': len(doc.paragraphs),
                'table_count': len(doc.tables)
            }

            return {
                'text': '\n'.join(text_content),
                'metadata': metadata,
                'success': True
            }

        except Exception as e:
            logger.error(f"DOCX parsing error: {e}")
            return {
                'text': '',
                'metadata': {},
                'success': False,
                'error': str(e)
            }

    def parse_pptx(self, file_path: str) -> Dict:
        """PowerPoint 문서 파싱"""
        try:
            prs = Presentation(file_path)
            text_content = []

            for slide_num, slide in enumerate(prs.slides):
                slide_text = [f"[Slide {slide_num + 1}]"]

                # 모든 도형의 텍스트 추출
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                text_content.append('\n'.join(slide_text))

            metadata = {
                'slide_count': len(prs.slides)
            }

            return {
                'text': '\n\n'.join(text_content),
                'metadata': metadata,
                'success': True
            }

        except Exception as e:
            logger.error(f"PPTX parsing error: {e}")
            return {
                'text': '',
                'metadata': {},
                'success': False,
                'error': str(e)
            }

    def parse_xlsx(self, file_path: str) -> Dict:
        """Excel 문서 파싱"""
        try:
            wb = load_workbook(file_path, data_only=True)
            text_content = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                sheet_text = [f"[Sheet: {sheet_name}]"]

                # 셀 데이터 추출
                for row in ws.iter_rows():
                    row_data = []
                    for cell in row:
                        if cell.value is not None:
                            row_data.append(str(cell.value))

                    if row_data:
                        sheet_text.append('\t'.join(row_data))

                text_content.append('\n'.join(sheet_text))

            metadata = {
                'sheet_count': len(wb.sheetnames),
                'sheet_names': wb.sheetnames
            }

            wb.close()

            return {
                'text': '\n\n'.join(text_content),
                'metadata': metadata,
                'success': True
            }

        except Exception as e:
            logger.error(f"XLSX parsing error: {e}")
            return {
                'text': '',
                'metadata': {},
                'success': False,
                'error': str(e)
            }

    def parse_image(self, file_path: str) -> Dict:
        """이미지 파일에서 OCR로 텍스트 추출"""
        try:
            if not PADDLE_AVAILABLE:
                return {
                    'text': '',
                    'metadata': {'ocr_available': False},
                    'success': True,
                    'error': 'PaddleOCR not installed'
                }

            # 이미지 읽기
            image = Image.open(file_path)

            # OCR 수행
            result = self.ocr_model.ocr(file_path, cls=True)

            # 텍스트 추출
            text_lines = []
            if result and result[0]:
                for line in result[0]:
                    text_lines.append(line[1][0])  # 인식된 텍스트

            metadata = {
                'image_size': image.size,
                'image_mode': image.mode,
                'ocr_confidence': 0.0
            }

            # 평균 신뢰도 계산
            if result and result[0]:
                confidences = [line[1][1] for line in result[0]]
                metadata['ocr_confidence'] = sum(confidences) / len(confidences) if confidences else 0.0

            image.close()

            return {
                'text': '\n'.join(text_lines),
                'metadata': metadata,
                'success': True
            }

        except Exception as e:
            logger.error(f"Image OCR error: {e}")
            return {
                'text': '',
                'metadata': {},
                'success': False,
                'error': str(e)
            }


def test_parser():
    """파서 테스트 함수"""
    parser = DocumentParser()

    test_files = [
        "/path/to/test.pdf",
        "/path/to/test.docx",
        "/path/to/test.pptx",
        "/path/to/test.xlsx",
        "/path/to/test.png",
    ]

    for file_path in test_files:
        if os.path.exists(file_path):
            print(f"\n{'='*60}")
            print(f"Testing: {file_path}")
            print('='*60)

            result = parser.parse_file(file_path)

            if result['success']:
                print(f"✅ Success!")
                print(f"Text length: {len(result['text'])} characters")
                print(f"Metadata: {result['metadata']}")
                print(f"\nFirst 200 characters:")
                print(result['text'][:200])
            else:
                print(f"❌ Failed: {result['error']}")


if __name__ == "__main__":
    test_parser()
