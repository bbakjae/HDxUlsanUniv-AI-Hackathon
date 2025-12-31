"""
합성 데이터 생성 스크립트
다양한 포맷의 문서 파일을 자동 생성하여 테스트 데이터셋 구축
"""

import os
import random
from datetime import datetime, timedelta
from pathlib import Path

# Office 문서 생성
from docx import Document
from docx.shared import Pt, Inches
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill

# PDF 생성
from reportlab.lib.pagesizes import letter, A4
from reportlab.pdfgen import canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib.units import inch

# 이미지 생성
from PIL import Image, ImageDraw, ImageFont

# 유틸리티
from tqdm import tqdm
import yaml


# 한국어 샘플 데이터
SAMPLE_DATA = {
    "departments": ["기획팀", "개발팀", "마케팅팀", "영업팀", "인사팀", "재무팀", "디자인팀", "품질관리팀"],
    "projects": [
        "신제품 개발 프로젝트",
        "고객 만족도 향상 계획",
        "AI 자동화 시스템 구축",
        "글로벌 시장 진출 전략",
        "조직 문화 개선 프로그램",
        "디지털 전환 로드맵",
        "공급망 최적화 방안",
        "ESG 경영 실천 계획"
    ],
    "topics": [
        "2024년 상반기 실적 보고",
        "차세대 기술 트렌드 분석",
        "경쟁사 벤치마킹 결과",
        "고객 이탈 방지 전략",
        "비용 절감 방안",
        "인재 채용 전략",
        "브랜드 가치 향상 방안",
        "데이터 분석 결과 요약"
    ],
    "content_templates": [
        "본 문서는 {topic}에 대한 상세한 분석 자료입니다.\n\n"
        "1. 현황 분석\n"
        "- 현재 상황: {department}에서 진행 중인 {project}의 진행 상황을 점검하였습니다.\n"
        "- 주요 성과: 목표 대비 85% 달성하였으며, 예상보다 빠른 진행을 보이고 있습니다.\n"
        "- 문제점: 일부 리소스 부족으로 인한 지연이 예상됩니다.\n\n"
        "2. 개선 방안\n"
        "- 단기: 추가 인력 투입 및 프로세스 간소화\n"
        "- 중기: 시스템 자동화를 통한 효율성 향상\n"
        "- 장기: 전사적 협업 체계 구축\n\n"
        "3. 기대 효과\n"
        "- 생산성 30% 향상\n"
        "- 비용 20% 절감\n"
        "- 고객 만족도 15% 증가",

        "{department}의 {project} 진행 보고서\n\n"
        "작성일: {date}\n"
        "작성자: {author}\n\n"
        "주요 내용:\n"
        "- {topic}에 대한 심층 분석을 완료하였습니다.\n"
        "- 시장 조사 결과, 향후 3년간 연평균 12% 성장이 예상됩니다.\n"
        "- 경쟁사 대비 우리의 강점은 기술력과 고객 서비스입니다.\n\n"
        "실행 계획:\n"
        "1분기: 시장 조사 및 기획\n"
        "2분기: 시스템 개발 및 테스트\n"
        "3분기: 파일럿 운영\n"
        "4분기: 전면 확대 적용\n\n"
        "예산: 약 5억원 소요 예상\n"
        "인력: 전담 인력 10명 배치 필요",

        "{topic} - {project} 최종 보고\n\n"
        "요약:\n"
        "본 보고서는 {department}에서 수행한 프로젝트의 최종 결과물입니다.\n\n"
        "성과 지표:\n"
        "- ROI: 150%\n"
        "- 목표 달성률: 92%\n"
        "- 고객 만족도: 4.5/5.0\n\n"
        "향후 계획:\n"
        "성공적인 결과를 바탕으로 다음 단계 프로젝트를 기획 중입니다.\n"
        "관련 부서와의 협업을 통해 시너지 효과를 극대화할 예정입니다.\n\n"
        "특이사항:\n"
        "프로젝트 수행 중 발견한 개선 기회를 별도 문서로 정리하였습니다."
    ],
    "authors": ["김철수", "이영희", "박민수", "정수진", "최동욱", "강미영", "윤서준", "임하늘"],
}


class SyntheticDataGenerator:
    def __init__(self, output_dir: str):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def generate_random_content(self) -> dict:
        """무작위 문서 내용 생성"""
        content_data = {
            "department": random.choice(SAMPLE_DATA["departments"]),
            "project": random.choice(SAMPLE_DATA["projects"]),
            "topic": random.choice(SAMPLE_DATA["topics"]),
            "author": random.choice(SAMPLE_DATA["authors"]),
            "date": (datetime.now() - timedelta(days=random.randint(0, 365))).strftime("%Y-%m-%d"),
        }

        template = random.choice(SAMPLE_DATA["content_templates"])
        content_data["content"] = template.format(**content_data)

        return content_data

    def generate_docx(self, filename: str):
        """Word 문서 생성"""
        doc = Document()
        data = self.generate_random_content()

        # 제목
        title = doc.add_heading(f"{data['topic']}", 0)

        # 메타 정보
        doc.add_paragraph(f"부서: {data['department']}")
        doc.add_paragraph(f"프로젝트: {data['project']}")
        doc.add_paragraph(f"작성자: {data['author']}")
        doc.add_paragraph(f"작성일: {data['date']}")
        doc.add_paragraph("")

        # 본문
        doc.add_heading("본문", level=1)
        doc.add_paragraph(data['content'])

        # 저장
        filepath = self.output_dir / filename
        doc.save(str(filepath))
        return filepath

    def generate_pptx(self, filename: str):
        """PowerPoint 문서 생성"""
        prs = Presentation()
        prs.slide_width = PptInches(10)
        prs.slide_height = PptInches(7.5)

        data = self.generate_random_content()

        # 제목 슬라이드
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = data['topic']
        subtitle.text = f"{data['department']} | {data['author']} | {data['date']}"

        # 내용 슬라이드
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = data['project']

        tf = body_shape.text_frame
        lines = data['content'].split('\n')
        for i, line in enumerate(lines[:10]):  # 최대 10줄
            if i == 0:
                tf.text = line
            else:
                p = tf.add_paragraph()
                p.text = line
                p.level = 1 if line.startswith('-') else 0

        # 저장
        filepath = self.output_dir / filename
        prs.save(str(filepath))
        return filepath

    def generate_xlsx(self, filename: str):
        """Excel 문서 생성"""
        wb = Workbook()
        ws = wb.active
        ws.title = "데이터 분석"

        data = self.generate_random_content()

        # 헤더
        ws['A1'] = "항목"
        ws['B1'] = "내용"
        ws['A1'].font = Font(bold=True, size=14)
        ws['B1'].font = Font(bold=True, size=14)

        # 메타 정보
        ws['A2'] = "부서"
        ws['B2'] = data['department']
        ws['A3'] = "프로젝트"
        ws['B3'] = data['project']
        ws['A4'] = "주제"
        ws['B4'] = data['topic']
        ws['A5'] = "작성자"
        ws['B5'] = data['author']
        ws['A6'] = "작성일"
        ws['B6'] = data['date']

        # 데이터 테이블
        ws['A8'] = "월"
        ws['B8'] = "목표"
        ws['C8'] = "실적"
        ws['D8'] = "달성률"

        for col in ['A8', 'B8', 'C8', 'D8']:
            ws[col].font = Font(bold=True)
            ws[col].fill = PatternFill(start_color="CCCCCC", end_color="CCCCCC", fill_type="solid")

        months = ["1월", "2월", "3월", "4월", "5월", "6월"]
        for i, month in enumerate(months):
            row = 9 + i
            target = random.randint(80, 120)
            actual = random.randint(70, 130)
            achievement = round((actual / target) * 100, 1)

            ws[f'A{row}'] = month
            ws[f'B{row}'] = target
            ws[f'C{row}'] = actual
            ws[f'D{row}'] = f"{achievement}%"

        # 열 너비 조정
        ws.column_dimensions['A'].width = 15
        ws.column_dimensions['B'].width = 30
        ws.column_dimensions['C'].width = 15
        ws.column_dimensions['D'].width = 15

        # 저장
        filepath = self.output_dir / filename
        wb.save(str(filepath))
        return filepath

    def generate_pdf(self, filename: str):
        """PDF 문서 생성 (한글 지원)"""
        filepath = self.output_dir / filename
        data = self.generate_random_content()

        # PDF 생성
        c = canvas.Canvas(str(filepath), pagesize=A4)
        width, height = A4

        # 제목 (영어로 대체 - 한글 폰트 없는 경우)
        c.setFont("Helvetica-Bold", 16)
        c.drawString(50, height - 50, "Document Report")

        # 메타 정보
        c.setFont("Helvetica", 12)
        y = height - 100
        c.drawString(50, y, f"Department: {data['department']}")
        c.drawString(50, y - 20, f"Project: {data['project']}")
        c.drawString(50, y - 40, f"Author: {data['author']}")
        c.drawString(50, y - 60, f"Date: {data['date']}")

        # 본문 (영어로 변환)
        c.setFont("Helvetica", 10)
        y = y - 100

        # 간단한 내용 작성
        lines = [
            "This is a synthetic document generated for testing purposes.",
            "",
            "Key Points:",
            "- Project status: On track",
            "- Budget: Within allocated range",
            "- Timeline: Meeting milestones",
            "- Quality: Exceeds expectations",
            "",
            "Next Steps:",
            "1. Continue monitoring progress",
            "2. Prepare for next phase",
            "3. Update stakeholders",
        ]

        for line in lines:
            c.drawString(50, y, line)
            y -= 15
            if y < 50:  # 새 페이지
                c.showPage()
                c.setFont("Helvetica", 10)
                y = height - 50

        c.save()
        return filepath

    def generate_image(self, filename: str):
        """이미지 생성 (텍스트 포함)"""
        # 이미지 생성
        img = Image.new('RGB', (800, 600), color=(255, 255, 255))
        d = ImageDraw.Draw(img)

        data = self.generate_random_content()

        # 배경색
        colors = [(230, 240, 250), (250, 240, 230), (240, 250, 230), (250, 230, 240)]
        bg_color = random.choice(colors)
        d.rectangle([(0, 0), (800, 600)], fill=bg_color)

        # 제목 박스
        d.rectangle([(50, 50), (750, 150)], fill=(100, 150, 200))

        # 텍스트 (기본 폰트 사용)
        try:
            # 시스템 폰트 시도
            font_large = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 24)
            font_medium = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 16)
        except:
            # 기본 폰트 사용
            font_large = ImageFont.load_default()
            font_medium = ImageFont.load_default()

        # 제목
        d.text((60, 80), data['topic'][:30], fill=(255, 255, 255), font=font_large)

        # 정보
        y = 200
        info_lines = [
            f"Department: {data['department']}",
            f"Project: {data['project']}",
            f"Author: {data['author']}",
            f"Date: {data['date']}",
        ]

        for line in info_lines:
            d.text((60, y), line, fill=(50, 50, 50), font=font_medium)
            y += 40

        # 도형 (데이터 시각화 흉내)
        for i in range(5):
            x = 60 + i * 140
            height_val = random.randint(50, 200)
            d.rectangle([(x, 550 - height_val), (x + 100, 550)],
                       fill=(random.randint(100, 200), random.randint(100, 200), random.randint(100, 200)))

        # 저장
        filepath = self.output_dir / filename
        img.save(str(filepath))
        return filepath

    def generate_dataset(self, num_files_per_type: int = 10):
        """전체 데이터셋 생성"""
        print(f"📁 합성 데이터 생성 시작: {self.output_dir}")
        print(f"각 파일 타입당 {num_files_per_type}개씩 생성합니다.\n")

        file_generators = [
            ("docx", self.generate_docx),
            ("pptx", self.generate_pptx),
            ("xlsx", self.generate_xlsx),
            ("pdf", self.generate_pdf),
            ("png", self.generate_image),
            ("jpg", self.generate_image),
        ]

        generated_files = []

        for file_type, generator_func in file_generators:
            print(f"\n🔧 {file_type.upper()} 파일 생성 중...")

            for i in tqdm(range(num_files_per_type)):
                # 파일명 생성 (한글 + 영어 혼합)
                dept = random.choice(SAMPLE_DATA["departments"])
                random_name = f"{dept}_{random.choice(['보고서', '계획서', '분석자료', '회의록'])}_{i+1:03d}"
                filename = f"{random_name}.{file_type}"

                try:
                    filepath = generator_func(filename)
                    generated_files.append(str(filepath))
                except Exception as e:
                    print(f"  ⚠️ {filename} 생성 실패: {e}")

        print(f"\n✅ 총 {len(generated_files)}개 파일 생성 완료!")
        print(f"📂 저장 위치: {self.output_dir}")

        # 통계 출력
        print("\n📊 생성된 파일 통계:")
        for file_type, _ in file_generators:
            count = len([f for f in generated_files if f.endswith(file_type)])
            print(f"  - {file_type.upper()}: {count}개")

        return generated_files


def main():
    # 설정 로드
    config_path = Path(__file__).parent.parent / "config" / "config.yaml"
    with open(config_path, 'r', encoding='utf-8') as f:
        config = yaml.safe_load(f)

    output_dir = config['data']['network_drive']

    # 생성기 초기화
    generator = SyntheticDataGenerator(output_dir)

    # 데이터셋 생성 (각 타입당 20개)
    generated_files = generator.generate_dataset(num_files_per_type=20)

    print("\n🎉 합성 데이터 생성이 완료되었습니다!")


if __name__ == "__main__":
    main()
